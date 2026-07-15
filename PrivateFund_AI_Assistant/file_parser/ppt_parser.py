"""PowerPoint 材料解析器。"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from models import SourceChunk


def parse_ppt(file_path: Path) -> list[SourceChunk]:
    """按幻灯片提取形状文字和表格文字。"""
    presentation = Presentation(str(file_path))
    chunks: list[SourceChunk] = []
    for page, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    texts.append(" | ".join(cell.text.strip() for cell in row.cells))
        text = "\n".join(dict.fromkeys(item for item in texts if item))
        if text:
            chunks.append(SourceChunk(text, file_path.name, f"第 {page} 页", "PPT"))
    return chunks

