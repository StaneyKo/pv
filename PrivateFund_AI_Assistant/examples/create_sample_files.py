"""生成可用于端到端测试的示例 Word、PPT、Excel、PDF 与 ZIP。"""

from __future__ import annotations

import zipfile
from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


EXAMPLE_DIR = Path(__file__).resolve().parent


def create_word() -> Path:
    """创建示例公司介绍 Word。"""
    path = EXAMPLE_DIR / "涵德投资_公司介绍.docx"
    document = Document()
    document.add_heading("涵德投资公司介绍", level=1)
    document.add_paragraph("涵德投资成立于2013年，专注于量化投资。本文件为软件功能测试材料。")
    document.add_heading("投研团队", level=2)
    document.add_paragraph("团队由投资经理、量化研究员和交易人员组成，具体人数资料未披露。")
    document.save(path)
    return path


def create_ppt() -> Path:
    """创建示例路演 PPT。"""
    path = EXAMPLE_DIR / "涵德投资_500指增路演.pptx"
    presentation = Presentation()
    for title, body in (
        ("中证500指数增强策略", "多因子选股结合组合优化，控制行业、市值及风格暴露。"),
        ("超额收益来源", "基本面、量价与分析师预期因子共同贡献选股收益。"),
        ("风险控制", "覆盖数据检查、模型检验、组合约束、交易监控和盘后复核。"),
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        box.text_frame.text = body
    presentation.save(path)
    return path


def create_excel() -> Path:
    """创建示例产品信息 Excel。"""
    path = EXAMPLE_DIR / "涵德投资_产品资料.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "产品信息"
    sheet.append(["产品名称", "策略", "管理规模", "最大回撤"])
    sheet.append(["示例一号", "中证500指数增强", "资料未披露", "资料未披露"])
    workbook.save(path)
    return path


def create_pdf() -> Path:
    """创建示例风险提示 PDF，使用英文以避免测试环境缺少中文字体。"""
    path = EXAMPLE_DIR / "sample_risk_note.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Risk Notice\nModel risk, market risk and liquidity risk should be reviewed.")
    document.save(path)
    document.close()
    return path


def create_zip(paths: list[Path]) -> Path:
    """将多类示例材料打包为 ZIP。"""
    path = EXAMPLE_DIR / "涵德投资_示例材料包.zip"
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for item in paths:
            archive.write(item, item.name)
    return path


def main() -> None:
    """生成全部示例文件。"""
    paths = [create_word(), create_ppt(), create_excel(), create_pdf()]
    create_zip(paths)
    print("示例文件已生成：")
    for path in paths + [EXAMPLE_DIR / "涵德投资_示例材料包.zip"]:
        print(path.name)


if __name__ == "__main__":
    main()

